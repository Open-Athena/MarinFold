"""Build the exp163 slide deck (python-pptx; no node in this environment)."""
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# --- palette: deep ocean + amber accent (contact maps / structure) -----------
MID   = C(0x21, 0x29, 0x5C)   # midnight  - dominant dark
TEAL  = C(0x1C, 0x72, 0x93)   # teal      - supporting
DEEP  = C(0x06, 0x5A, 0x82)   # deep blue - supporting
AMBER = C(0xE8, 0x8B, 0x2F)   # amber     - sharp accent (key numbers)
CORAL = C(0xC4, 0x45, 0x36)   # coral     - failure / negative
WHITE = C(0xFF, 0xFF, 0xFF)
INK   = C(0x1A, 0x1F, 0x36)
MUTE  = C(0x6B, 0x74, 0x8C)
CARD  = C(0xF1, 0xF5, 0xF9)
GREEN = C(0x2E, 0x7D, 0x5B)

HEAD = "Cambria"      # safe-list serif header
BODY = "Calibri"      # safe-list sans body
MONO = "Courier New"  # safe-list mono

prs = Presentation()
prs.slide_width, prs.slide_height = In(13.333), In(7.5)
W, H = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def sl(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = MID if dark else WHITE
    return s


def box(s, x, y, w, h, text, size=16, color=INK, bold=False, font=BODY,
        align=PP_ALIGN.LEFT, italic=False, spacing=1.0, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = ln
        f = r.font
        f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, font
        f.color.rgb = color
    return tb


def bullets(s, x, y, w, h, items, size=15, color=INK, gap=8, lead=None):
    tb = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.15
        if isinstance(it, tuple):
            head, rest = it          # NB: not `lead` -- that is the colour param
            r = p.add_run(); r.text = head
            r.font.size, r.font.bold, r.font.name = Pt(size), True, BODY
            # MID on a dark slide is invisible; callers on dark backgrounds pass lead=
            r.font.color.rgb = lead if lead is not None else MID
            r2 = p.add_run(); r2.text = rest
            r2.font.size, r2.font.name = Pt(size), BODY
            r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = it
            r.font.size, r.font.name = Pt(size), BODY
            r.font.color.rgb = color
    return tb


def card(s, x, y, w, h, fill=CARD, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(x), In(y), In(w), In(h))
    sh.adjustments[0] = 0.06
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(1.25)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def circle_num(s, x, y, d, n, fill=TEAL, tcol=WHITE, size=17):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, In(x), In(y), In(d), In(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    tf = sh.text_frame; tf.text = ""
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.size, r.font.bold, r.font.name = Pt(size), True, BODY
    r.font.color.rgb = tcol
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


def title(s, t, sub=None, dark=False):
    # One line only: a wrapped title would run into the subtitle beneath it.
    size = 34 if len(t) <= 40 else (30 if len(t) <= 50 else 26)
    box(s, 0.7, 0.5, W - 1.4, 0.85, t, size=size, bold=True, font=HEAD,
        color=WHITE if dark else MID)
    if sub:
        box(s, 0.7, 1.32, W - 1.4, 0.5, sub, size=15,
            color=C(0xC9, 0xD4, 0xE8) if dark else MUTE, italic=True)


def stat(s, x, y, w, value, label, vcol=AMBER, vsize=44, sub=None):
    box(s, x, y, w, 0.75, value, size=vsize, bold=True, font=HEAD, color=vcol,
        align=PP_ALIGN.CENTER)
    box(s, x, y + 0.78, w, 0.55, label, size=13, color=MID, align=PP_ALIGN.CENTER,
        bold=True)
    if sub:
        box(s, x, y + 1.22, w, 0.5, sub, size=11, color=MUTE, align=PP_ALIGN.CENTER)


def mono(s, x, y, w, h, lines, size=10.5, fill=C(0x14, 0x1A, 0x33), color=C(0xDC, 0xE6, 0xF5)):
    card(s, x, y, w, h, fill=fill)
    tb = s.shapes.add_textbox(In(x + 0.22), In(y + 0.16), In(w - 0.44), In(h - 0.32))
    tf = tb.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.06
        p.alignment = PP_ALIGN.LEFT
        txt, col, bold = (ln if isinstance(ln, tuple) else (ln, color, False))
        r = p.add_run(); r.text = txt
        r.font.size, r.font.name, r.font.bold = Pt(size), MONO, bold
        r.font.color.rgb = col
    return tb


def table(s, x, y, w, rows, colw, size=12.5, hdr_fill=MID, row_h=0.36, hdr_h=0.42,
          bold_col0=False, hilite=None):
    """rows[0] is the header. hilite = set of (row_idx, col_idx) -> amber bold."""
    n = len(rows)
    tbl_h = hdr_h + row_h * (n - 1)
    shp = s.shapes.add_table(n, len(rows[0]), In(x), In(y), In(w), In(tbl_h)).table
    for j, cw in enumerate(colw):
        shp.columns[j].width = Emu(int(In(cw)))
    shp.rows[0].height = Emu(int(In(hdr_h)))
    for i in range(1, n):
        shp.rows[i].height = Emu(int(In(row_h)))
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = shp.cell(i, j)
            c.margin_left = In(0.10); c.margin_right = In(0.08)
            c.margin_top = In(0.03); c.margin_bottom = In(0.03)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = hdr_fill if i == 0 else (WHITE if i % 2 else CARD)
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            f = r.font
            f.size = Pt(size); f.name = BODY
            hot = hilite and (i, j) in hilite
            f.bold = (i == 0) or hot or (bold_col0 and j == 0)
            f.color.rgb = WHITE if i == 0 else (AMBER if hot else INK)
    return shp


def chart_style(ch, legend=False):
    ch.has_title = False
    ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(11); ch.legend.font.name = BODY
    for ax in (ch.category_axis, ch.value_axis):
        ax.tick_labels.font.size = Pt(11)
        ax.tick_labels.font.name = BODY
        ax.tick_labels.font.color.rgb = MUTE
    ch.value_axis.has_major_gridlines = True
    gl = ch.value_axis.major_gridlines.format.line
    gl.color.rgb = C(0xDD, 0xE3, 0xEC); gl.width = Pt(0.75)
    ch.category_axis.has_major_gridlines = False


def footer(s, txt, dark=False):
    box(s, 0.7, H - 0.62, W - 1.4, 0.35, txt, size=10.5,
        color=C(0x8F, 0x9C, 0xB8) if dark else MUTE, italic=True)


# ============================================================ 1. title
s = sl(dark=True)
box(s, 0.9, 2.15, 11.5, 1.1, "Teaching contacts-v1 to refine\nits own candidate rollouts",
    size=40, bold=True, font=HEAD, color=WHITE, spacing=1.12)
box(s, 0.9, 3.75, 10.5, 0.5,
    "exp163  ·  can K noisy samples be aggregated in-context into a better contact set?",
    size=16, color=C(0x9F, 0xC4, 0xE0))
ln = s.shapes.add_shape(MSO_SHAPE.OVAL, In(0.9), In(4.55), In(0.16), In(0.16))
ln.fill.solid(); ln.fill.fore_color.rgb = AMBER; ln.line.fill.background(); ln.shadow.inherit = False
box(s, 1.22, 4.5, 10.5, 0.4,
    "Phases 0–2 complete · Phase 3 in flight · MarinFold issue #163 / PR #164",
    size=13, color=C(0xC9, 0xD4, 0xE8))
box(s, 0.9, 6.5, 10.5, 0.4, "2026-07-27", size=12, color=C(0x7F, 0x92, 0xB5))

# ============================================================ 2. problem
s = sl()
title(s, "The information is there — spread across samples",
      "A single contacts-v1 rollout is poor. Sixteen of them collectively are not.")
for i, (v, l, sub, col) in enumerate([
        ("0.122", "single rollout", "mean R-precision", MUTE),
        ("0.233", "oracle best-of-16", "not deployable", TEAL),
        ("0.516", "union of 16", "recall ceiling", AMBER)]):
    x = 0.9 + i * 4.0
    card(s, x, 2.15, 3.6, 2.05)
    stat(s, x, 2.4, 3.6, v, l, vcol=col, sub=sub)
box(s, 0.9, 4.6, 11.5, 1.5,
    "So: can the model be taught to read K of its own rollouts and write out a contact set\n"
    "better than any of them — an aggregation operator learned in context?",
    size=17, color=INK, spacing=1.25)
card(s, 0.9, 5.75, 11.5, 0.95, fill=C(0xFD, 0xF3, 0xE3))
box(s, 1.15, 5.98, 11.0, 0.6,
    "The bar is not “beat one rollout”. It is: beat training-free consensus voting, "
    "and beat the model’s own one-shot prediction.",
    size=14, color=C(0x8A, 0x53, 0x12), bold=True)

# ============================================================ 3. probe 1
s = sl()
title(s, "Probe 1 — aggregation beats selection",
      "exp98: 1,000 proteins × 1,000 rollouts. R-precision, all band.")
cd = CategoryChartData()
cd.categories = ["K = 8", "K = 16", "K = 32"]
cd.add_series("mean single", (0.122, 0.122, 0.122))
cd.add_series("proxy best-of-K (nll)", (0.137, 0.145, 0.144))
cd.add_series("consensus vote", (0.194, 0.224, 0.244))
cd.add_series("oracle best-of-K", (0.208, 0.233, 0.252))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, In(0.9), In(2.05), In(7.3), In(4.0), cd)
ch = gf.chart
chart_style(ch, legend=True)
for ser, col in zip(ch.plots[0].series, [MUTE, C(0x9A, 0xA7, 0xBD), TEAL, MID]):
    ser.format.fill.solid(); ser.format.fill.fore_color.rgb = col
ch.value_axis.maximum_scale = 0.28
ch.value_axis.minimum_scale = 0
bullets(s, 8.55, 2.2, 3.9, 3.8, [
    ("Voting ≈ oracle. ", "Consensus reaches ~97% of the non-deployable oracle "
     "selection, and keeps rising with K."),
    ("Confidence is useless. ", "Within-protein Spearman(nll, F1) median −0.17; the "
     "min-nll pick sits at F1 percentile 0.525 — random."),
    ("⇒ No deployable way to rank or filter candidates.", ""),
], size=14, gap=13)
card(s, 8.55, 5.5, 3.9, 0.85, fill=MID)
box(s, 8.8, 5.72, 3.5, 0.5, "So: train on unordered candidate sets.",
    size=14, color=WHITE, bold=True)

# ============================================================ 4. probe 2
s = sl()
title(s, "Probe 2 — the control that reframes everything",
      "Paired on the same 149 proteins: does voting actually beat the base model?")
table(s, 0.9, 2.15, 6.6, [
    ["method", "R-prec", "AUC"],
    ["base E8 calibrated matrix", "0.221", "0.887"],
    ["consensus vote (K=16)", "0.224", "0.735"],
    ["single rollout", "0.125", "—"],
    ["oracle best-of-16", "0.241", "—"],
    ["union-recall ceiling", "0.530", "—"],
], [3.5, 1.6, 1.5], hilite={(1, 1), (1, 2), (5, 1)})
card(s, 7.9, 2.15, 4.55, 2.3, fill=MID)
box(s, 8.2, 2.45, 3.95, 0.5, "voting − matrix", size=14, color=C(0x9F, 0xC4, 0xE0))
box(s, 8.2, 2.9, 3.95, 0.9, "+0.003", size=44, bold=True, font=HEAD, color=AMBER)
box(s, 8.2, 3.75, 3.95, 0.5, "vote wins on 50% of proteins — a tie",
    size=13, color=WHITE)
bullets(s, 0.9, 4.85, 11.5, 1.9, [
    ("Consensus voting is a Monte-Carlo estimate of the model's own per-pair marginal. ",
     "It adds nothing the base model does not already emit."),
    ("⇒ A refiner that learns only consensus is worthless. ",
     "It has to extract supra-marginal signal — joint structure the marginal cannot express."),
], size=15, gap=12)

# ============================================================ 5. probe 3 (crux)
s = sl()
title(s, "Probe 3 — the crux",
      "Zero-shot, paired cond-vs-uncond on the identical remaining-pair task. Joint structure exists — but it is precision-gated.")
cd = CategoryChartData()
cd.categories = ["oracle: 50% of TRUE contacts", "noisy: one real rollout (~13% prec)"]
cd.add_series("unconditional", (0.145, 0.179))
cd.add_series("conditioned", (0.556, 0.092))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, In(0.9), In(2.05), In(7.1), In(3.9), cd)
ch = gf.chart
chart_style(ch, legend=True)
ch.plots[0].series[0].format.fill.solid()
ch.plots[0].series[0].format.fill.fore_color.rgb = MUTE
ch.plots[0].series[1].format.fill.solid()
ch.plots[0].series[1].format.fill.fore_color.rgb = TEAL
pl = ch.plots[0]
pl.has_data_labels = True
pl.data_labels.number_format = "0.000"
pl.data_labels.number_format_is_linked = False
pl.data_labels.font.size = Pt(11); pl.data_labels.font.name = BODY
pl.data_labels.font.bold = True
pl.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
card(s, 8.35, 2.15, 4.1, 1.75, fill=C(0xE8, 0xF3, 0xEE))
box(s, 8.6, 2.38, 3.6, 0.45, "TRUE partial context", size=13, bold=True, color=GREEN)
box(s, 8.6, 2.82, 3.6, 0.75, "+0.41", size=34, bold=True, font=HEAD, color=GREEN)
box(s, 8.6, 3.45, 3.6, 0.4, "AUC → 0.99 · better on 100%", size=11.5, color=C(0x2E, 0x5D, 0x4A))
card(s, 8.35, 4.05, 4.1, 1.75, fill=C(0xFA, 0xEC, 0xE9))
box(s, 8.6, 4.28, 3.6, 0.45, "NOISY candidate context", size=13, bold=True, color=CORAL)
box(s, 8.6, 4.72, 3.6, 0.75, "−0.087", size=34, bold=True, font=HEAD, color=CORAL)
box(s, 8.6, 5.35, 3.6, 0.4, "worse on 91% of proteins", size=11.5, color=C(0x8E, 0x3A, 0x2C))
card(s, 0.9, 6.15, 11.55, 0.85, fill=MID)
box(s, 1.15, 6.36, 11.0, 0.5,
    "The base model trusts its context as ground truth — so training, not prompting, is what's required.",
    size=14.5, color=WHITE, bold=True)

# ============================================================ 6. format
s = sl()
title(s, "The document format — zero vocab change",
      "A normal contacts-v1 document with K candidate blocks spliced in.")
mono(s, 0.9, 2.05, 7.5, 2.35, [
    ("<contacts-v1> <begin_sequence> …sequence…", C(0x9A, 0xB4, 0xD8), False),
    ("  <CAND> <contact> pi pj …        ← candidate 1", C(0xE8, 0x8B, 0x2F), False),
    ("  …                                  K blocks, unordered", C(0x77, 0x88, 0xAA), False),
    ("  <CAND> <contact> pi pj …        ← candidate K", C(0xE8, 0x8B, 0x2F), False),
    ("<begin_statements> …TRUE contacts… <end>", C(0x7A, 0xD6, 0xA8), True),
], size=12.5)
box(s, 0.9, 4.55, 7.5, 0.45, "Loss is armed on the green span only.", size=13,
    color=MUTE, italic=True)
bullets(s, 8.75, 2.1, 3.7, 4.4, [
    ("<CAND> is a spare token. ", "It reuses <contacts-and-distances-v1> — the other "
     "format's sentinel, never emitted inside a contacts-v1 doc. Weight-compatible "
     "with every checkpoint."),
    ("K separate blocks. ", "Recurrence across blocks is a precision cue the model "
     "can read. Merging would destroy it."),
    ("K ~ Uniform{0…16}. ", "K=0 is a plain document — guards against both ignoring "
     "candidates and blindly trusting them."),
], size=13, gap=11)
footer(s, "Per-candidate subsampling uniform[1, n_pred] keeps test-time full candidates in distribution.")

# ============================================================ 7. real example
s = sl()
title(s, "A real training document (K = 5)",
      "251 tokens — the shortest in the corpus, shown near-complete. Typical is 2,270 tokens / ~424 candidate contacts.")
mono(s, 0.9, 2.0, 11.55, 3.35, [
    ("# sequence header — loss weight 0", C(0x77, 0x88, 0xAA), False),
    ("<contacts-v1> <begin_sequence> <p16> <ALA> <p25> <LEU> <p12> <ALA> <p1973> <SER> …", C(0x9A, 0xB4, 0xD8), False),
    ("", C(0x77, 0x88, 0xAA), False),
    ("# candidate block 1/5  (1 contact)     — loss weight 0", C(0x77, 0x88, 0xAA), False),
    ("<contacts-and-distances-v1> <contact> <p1986> <p1992>", C(0xE8, 0x8B, 0x2F), False),
    ("# candidate block 2/5  (6 contacts)    — loss weight 0", C(0x77, 0x88, 0xAA), False),
    ("<contacts-and-distances-v1> <contact> <p1988> <p1970> <contact> <p1993> <p1976> …", C(0xE8, 0x8B, 0x2F), False),
    ("# candidate blocks 3–5  (9, 2, 11 contacts) — loss weight 0", C(0x77, 0x88, 0xAA), False),
    ("", C(0x77, 0x88, 0xAA), False),
    ("# answer span — loss weight 1  (9 TRUE contacts + <end>)", C(0x77, 0x88, 0xAA), False),
    ("<begin_statements> <contact> <p1> <p1993> <contact> <p1980> <p1987>", C(0x7A, 0xD6, 0xA8), True),
    ("<contact> <p1980> <p1967> <contact> <p14> <p7> … <end>", C(0x7A, 0xD6, 0xA8), True),
], size=11.5)
card(s, 0.9, 5.6, 11.55, 0.8, fill=C(0xFD, 0xF3, 0xE3))
box(s, 1.15, 5.8, 11.0, 0.45,
    "29 distinct contacts proposed across the five blocks · 0 of them appear in the answer. "
    "That is the regime the refiner must survive.",
    size=13.5, color=C(0x8A, 0x53, 0x12), bold=True)

# ============================================================ 8. recurrence
s = sl()
title(s, "The format's central bet, measured",
      "Does recurrence across candidate blocks actually carry precision? 1,393 sampled documents.")
cd = CategoryChartData()
cd.categories = ["appears in\nexactly 1 block", "all distinct\ncandidates", "appears in\n>1 block"]
cd.add_series("precision vs the true answer", (0.079, 0.097, 0.247))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, In(0.9), In(2.15), In(7.0), In(4.05), cd)
ch = gf.chart
chart_style(ch, legend=False)
pts = ch.plots[0]
pts.has_data_labels = True
pts.data_labels.number_format = "0.000"
pts.data_labels.number_format_is_linked = False
pts.data_labels.font.size = Pt(13); pts.data_labels.font.bold = True
pts.data_labels.font.name = BODY
pts.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
ser = ch.plots[0].series[0]
for idx, col in enumerate([MUTE, TEAL, AMBER]):
    ser.points[idx].format.fill.solid()
    ser.points[idx].format.fill.fore_color.rgb = col
ch.value_axis.maximum_scale = 0.30
ch.value_axis.minimum_scale = 0
card(s, 8.25, 2.3, 4.2, 2.15, fill=MID)
box(s, 8.55, 2.55, 3.7, 0.5, "a recurring contact is", size=13, color=C(0x9F, 0xC4, 0xE0))
box(s, 8.55, 2.98, 3.7, 0.85, "3× ", size=44, bold=True, font=HEAD, color=AMBER)
box(s, 8.55, 3.75, 3.7, 0.5, "more likely to be true\nthan a singleton", size=13, color=WHITE)
bullets(s, 8.25, 4.75, 4.2, 1.9, [
    "The signal is present in the raw input, before any training — the model only has to learn to use it.",
    "It also explains why an explicit consensus block bought nothing over raw blocks in the MVP.",
], size=12.5, gap=10)

# ============================================================ 9. loss mask
s = sl()
title(s, "The loss mask, and where it had to move",
      "Armed on [<begin_statements>, <end>) — trains the contacts and the stop token.")
mono(s, 0.9, 2.0, 6.5, 1.5, [
    ("opened[i] = # <begin_statements> in ids[:i+1]", C(0x9A, 0xB4, 0xD8), False),
    ("closed[i] = # <end>              in ids[:i+1]", C(0x9A, 0xB4, 0xD8), False),
    ("weight[i] = 1.0 if opened[i] > closed[i]", C(0x7A, 0xD6, 0xA8), True),
], size=12)
box(s, 0.9, 3.62, 6.5, 0.45, "Two running counts ⇒ packing-safe: re-arms per document.",
    size=12.5, color=MUTE, italic=True)
mono(s, 0.9, 4.15, 6.5, 2.25, [
    ("i=453 <TYR>              w=0.0 → <begin_statements>", C(0x9A, 0xB4, 0xD8), False),
    ("i=454 <begin_statements> w=1.0 → <contact>", C(0x7A, 0xD6, 0xA8), True),
    ("i=455 <contact>          w=1.0 → <p133>", C(0x7A, 0xD6, 0xA8), True),
    ("  …", C(0x77, 0x88, 0xAA), False),
    ("i=976 <p145>             w=1.0 → <end>", C(0x7A, 0xD6, 0xA8), True),
    ("i=977 <end>              w=0.0 → <eos>", C(0x9A, 0xB4, 0xD8), False),
    ("i=978 <eos>              w=0.0 → next document", C(0x9A, 0xB4, 0xD8), False),
], size=11)
card(s, 7.75, 2.0, 4.7, 4.4, fill=CARD)
box(s, 8.05, 2.25, 4.1, 0.45, "levanter 1.2 removed the hook", size=15, bold=True, color=MID)
bullets(s, 8.05, 2.8, 4.1, 3.4, [
    "DatasetComponent.loss_weight_fn no longer exists. Per-token weights now reach training only through the cache, via PrebuiltLmDatasetFormat.",
    "So the mask is materialized offline: tokenize → greedy-pack to 8192 → write input_ids + loss_weights.",
    ("Strictly better. ", "Token ids are resolved from the live tokenizer, and nothing "
     "cloudpickles by module reference — so the GPU worker never imports experiment code."),
], size=12, gap=10)

# ============================================================ 10. MVP
s = sl()
title(s, "MVP — mechanism confirmed, margin overfitting-limited",
      "Local LoRA on 900 proteins, 60 held out. R-precision (all / long).")
table(s, 0.9, 2.15, 11.55, [
    ["model", "K = 0", "raw K = 16", "consensus block (44% prec)"],
    ["base E8", "0.229 / 0.153", "0.017 / 0.016", "0.149 / 0.124"],
    ["refiner v1 (r16)", "0.213 / 0.129", "0.244 / 0.169", "0.247 / 0.163"],
    ["refiner v2 (r32)", "0.183 / 0.120", "0.218 / 0.156", "0.214 / 0.146"],
], [3.0, 2.6, 2.75, 3.2], size=13, hilite={(1, 2), (2, 2)})
bullets(s, 0.9, 4.2, 5.6, 2.6, [
    ("Conditioning is learned. ", "+0.03 all / +0.04 long over the same model's K=0, "
     "stable across v1/v2 and reproduced in the sampling regime."),
    ("Base E8 is poisoned. ", "0.229 → 0.017 on identical input — the cleanest "
     "demonstration that training is what makes candidates usable."),
], size=13, gap=11)
bullets(s, 6.85, 4.2, 5.6, 2.6, [
    ("Preprocessing is a dead lever. ", "A 44%-precision consensus block ≈ raw blocks; "
     "the refiner already internalizes consensus."),
    ("Overfitting caps the margin. ", "v2 fit harder and got worse held-out "
     "(0.213 → 0.183). Scale is the lever — which motivated everything that follows."),
], size=13, gap=11)

# ============================================================ 11. generation
s = sl()
title(s, "Scale — generation validated on CoreWeave",
      "ESM-Atlas proteins, 16 × 1×H100 at batch priority, ~1h wall.")
for i, (v, l, sub, col) in enumerate([
        ("225,072", "rollouts", "over 9,375 proteins", TEAL),
        ("201.0", "mean n_pred", "vs GT ~199 — top-k fix confirmed", AMBER),
        ("1.00", "frac finished", "no truncation", GREEN),
        ("6,569", "packed sequences", "44.3M tokens, 82.3% density", MID)]):
    x = 0.75 + i * 3.05
    card(s, x, 2.15, 2.8, 2.15)
    stat(s, x, 2.4, 2.8, v, l, vcol=col, vsize=30, sub=sub)
bullets(s, 0.9, 4.7, 11.5, 2.1, [
    ("--top-k -1 is the fix. ", "exp98's top-k=50 rollouts averaged ~95 contacts against "
     "a GT of ~199; disabling top-k restores the length distribution (#142)."),
    ("Preemption is free. ", "One of 16 shards was preempted — the worker is resume-safe, "
     "which is exactly why 9,375 ≠ 10,000."),
    ("Corpus: ", "18,750 documents (2 per protein, K ~ U{0..16}), 0 OOV, 0 over-budget."),
], size=14, gap=11)

# ============================================================ 12. training
s = sl()
title(s, "Scale — training (Phase 2, complete)",
      "Warm-started from E8, 1-epoch cosine, batch 128, seq 8192, 8×H100 batch band, ~25 min/arm.")
table(s, 0.9, 2.15, 11.55, [
    ["peak LR", "train/loss (masked objective)", "base-task eval loss", "base-task bpb"],
    ["1e-4", "3.985 → 2.3979", "3.16941", "0.39489"],
    ["3e-4", "3.833 → 2.3915", "3.40526", "0.42428"],
], [2.4, 4.15, 2.7, 2.3], size=14, row_h=0.44, hilite={(1, 3)})
card(s, 0.9, 3.85, 5.6, 1.5, fill=C(0xE8, 0xF3, 0xEE))
box(s, 1.2, 4.08, 5.0, 1.05,
    "1e-4 wins. 3e-4 fits the objective 0.3% better but degrades base-task retention 7.4% more.",
    size=14, color=C(0x22, 0x50, 0x3E), bold=True)
card(s, 6.85, 3.85, 5.6, 1.5, fill=CARD)
box(s, 7.15, 4.08, 5.0, 1.05,
    "Re-running the 1e-4 arm reproduced an earlier run bit-identically — a free determinism "
    "check over tokenize → pack → cache → train.",
    size=14, color=INK)
bullets(s, 0.9, 5.6, 11.5, 1.3, [
    ("Same fitting-vs-retention tension as the MVP. ", "Both eval curves dip at step 13 "
     "then recover — LR-warmup damage, not divergence."),
    ("The HF export is free. ", "Each run writes checkpoints/<run>/hf/step-N with the "
     "tokenizer co-located, so Phase 3 needs no separate export step."),
], size=14, gap=10)

# ============================================================ 13. measurement trap
s = sl()
title(s, "A measurement trap worth carrying forward",
      "The planned warm-start check — “step-0 val loss ≈ 2.7566” — does not work.")
table(s, 0.9, 2.2, 9.6, [
    ["", "loss", "bpb", "implied bytes/token"],
    ["Eric's E8 (step 35679)", "2.75660", "0.39151", "10.16"],
    ["this harness (step 51)", "3.16941", "0.39489", "11.58"],
], [3.8, 2.0, 1.9, 1.9], size=14, row_h=0.44, hilite={(1, 2), (2, 2)})
card(s, 0.9, 4.0, 11.55, 0.8, fill=MID)
box(s, 1.15, 4.2, 11.0, 0.5,
    "loss ratio 1.1497  =  bpb ratio 1.0086  ×  bytes/token ratio 1.1399   — exactly.",
    size=15, color=WHITE, bold=True, font=MONO)
bullets(s, 0.9, 5.1, 11.5, 1.8, [
    ("No step-0 eval exists. ", "levanter fires eval hooks at multiples of steps_per_eval, "
     "so the first recorded value already carries warmup damage."),
    ("Per-token loss is not comparable across harnesses. ", "bpb divides by per-token-type "
     "byte lengths weighted by the loss mask, so packing changes the loss scale at "
     "identical model quality. bpb agrees to 0.9% — that is the evidence E8 loaded."),
    ("⚠ #137, #150 and #155 all quote cross-harness per-token loss targets.", ""),
], size=13.5, gap=10)

# ============================================================ 14. phase 3 results
s = sl()
title(s, "Phase 3 — the verdict", "553 exp89 eval proteins, paired, identical candidate contexts for both models.")
table(s, 0.75, 2.0, 11.85, [
    ["model", "K0", "K1", "K2", "K4", "K8", "K16", "consensus"],
    ["base", "0.3355", "0.1452", "0.1024", "0.0665", "0.0269", "0.0194", "0.2023"],
    ["refiner", "0.1978", "0.1555", "0.0813", "0.0383", "0.0211", "0.0165", "0.2220"],
    ["candidate contacts", "0", "98", "195", "391", "780", "1557", "58"],
], [2.55, 1.35, 1.2, 1.2, 1.2, 1.2, 1.2, 1.55], size=12.5, row_h=0.42,
   hilite={(1, 1), (2, 7)})
card(s, 0.75, 3.9, 5.75, 1.25, fill=C(0xFA, 0xEC, 0xE9))
box(s, 1.05, 4.1, 5.15, 0.85,
    "Kill criterion MET. The base model's one-shot matrix (0.3355) is still the best "
    "deployable prediction; the refiner's best arm is 0.2220.",
    size=13.5, color=C(0x8E, 0x3A, 0x2C), bold=True)
card(s, 6.85, 3.9, 5.75, 1.25, fill=C(0xE8, 0xF3, 0xEE))
box(s, 7.15, 4.1, 5.15, 0.85,
    "But raw blocks hurt BOTH models monotonically — even at K=1, comfortably "
    "in-distribution. Only the precision-filtered consensus block helps.",
    size=13.5, color=C(0x22, 0x50, 0x3E), bold=True)
box(s, 0.75, 5.45, 11.85, 0.45, "Harness check: base K0 = 0.3355 vs exp89's published E8 value of 0.3389.",
    size=13, color=MUTE, italic=True)
footer(s, "R-precision, all band. Long band tells the same story (base K0 0.2697, refiner consensus 0.1767).")

# ============================================================ 15. mechanism works
s = sl()
title(s, "The mechanism does work", "Same 56.7%-precision consensus block, fed to both models.")
cd = CategoryChartData()
cd.categories = ["base", "refiner"]
cd.add_series("K0 (no candidates)", (0.3355, 0.1978))
cd.add_series("+ consensus block", (0.2023, 0.2220))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, In(0.9), In(2.1), In(7.0), In(3.9), cd)
ch = gf.chart
chart_style(ch, legend=True)
ch.plots[0].series[0].format.fill.solid()
ch.plots[0].series[0].format.fill.fore_color.rgb = MUTE
ch.plots[0].series[1].format.fill.solid()
ch.plots[0].series[1].format.fill.fore_color.rgb = TEAL
pl = ch.plots[0]
pl.has_data_labels = True
pl.data_labels.number_format = "0.000"
pl.data_labels.number_format_is_linked = False
pl.data_labels.font.size = Pt(11); pl.data_labels.font.name = BODY
pl.data_labels.font.bold = True
pl.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
ch.value_axis.maximum_scale = 0.40
ch.value_axis.minimum_scale = 0
card(s, 8.3, 2.15, 4.15, 1.7, fill=C(0xFA, 0xEC, 0xE9))
box(s, 8.6, 2.36, 3.6, 0.42, "base: POISONED", size=13, bold=True, color=CORAL)
box(s, 8.6, 2.78, 3.6, 0.7, "-0.133", size=32, bold=True, font=HEAD, color=CORAL)
box(s, 8.6, 3.38, 3.6, 0.4, "gains on only 8% of proteins", size=11.5, color=C(0x8E, 0x3A, 0x2C))
card(s, 8.3, 4.0, 4.15, 1.7, fill=C(0xE8, 0xF3, 0xEE))
box(s, 8.6, 4.21, 3.6, 0.42, "refiner: USES IT", size=13, bold=True, color=GREEN)
box(s, 8.6, 4.63, 3.6, 0.7, "+0.024", size=32, bold=True, font=HEAD, color=GREEN)
box(s, 8.6, 5.23, 3.6, 0.4, "gains on 63% of proteins", size=11.5, color=C(0x2E, 0x5D, 0x4A))
card(s, 0.9, 6.2, 11.55, 0.8, fill=MID)
box(s, 1.15, 6.4, 11.0, 0.5,
    "A +0.157 swing in how the two models respond to identical input. Candidate discrimination WAS learned.",
    size=14.5, color=WHITE, bold=True)

# ============================================================ 16. why it loses
s = sl()
title(s, "Why it still loses: forgetting", "The refiner is not bad at refining. It is bad at contacts.")
for i, (v, l, sub, col) in enumerate([
        ("0.3355", "base K0", "one-shot contact ability", TEAL),
        ("0.1978", "refiner K0", "after 1 epoch of fine-tuning", CORAL),
        ("-41%", "of the base task", "base wins on 94% of proteins", CORAL)]):
    x = 0.9 + i * 4.0
    card(s, x, 2.15, 3.6, 2.0)
    stat(s, x, 2.4, 3.6, v, l, vcol=col, vsize=32, sub=sub)
bullets(s, 0.9, 4.5, 11.5, 2.4, [
    ("A +0.024 conditioning gain cannot climb out of a 0.138 hole. ",
     "Nothing clears the bar because the refiner starts from a much weaker model."),
    ("The LM proxy badly under-reports the damage. ",
     "Base-task val bpb moved 0.9% (0.39151 -> 0.39489). The contact metric moved 41%. "
     "Exactly what #89's loss-vs-R-precision study warned about."),
    ("Scaling 100x does not fix this. ",
     "Fix the forgetting first: LoRA (as the MVP used), lower LR, or mix plain "
     "contacts-v1 documents in so K=0 stays anchored."),
], size=14, gap=13)

# ============================================================ 17. v2 format
s = sl()
title(s, "v2 — one marker, for RL", "<begin_statements> means \u201cdiscard the previous structure; here is a new candidate\u201d.")
mono(s, 0.9, 2.05, 7.4, 2.5, [
    ("<contacts-v1> <begin_sequence> \u2026seq\u2026", C(0x9A, 0xB4, 0xD8), False),
    ("<begin_statements> \u2026draft 1\u2026            (no <end>)", C(0xE8, 0x8B, 0x2F), False),
    ("<begin_statements> \u2026draft 2\u2026", C(0xE8, 0x8B, 0x2F), False),
    ("<begin_statements> \u2026TRUE contacts\u2026 <end>", C(0x7A, 0xD6, 0xA8), True),
], size=12)
box(s, 0.9, 4.7, 7.4, 0.55, "Drafts and the answer are the SAME object — so one generation can emit N structures and be RL'd on best-of-N.", size=13, color=MUTE, italic=True, spacing=1.2)
bullets(s, 8.6, 2.1, 3.85, 4.4, [
    ("Only the final section closes. ", "Drafts are superseded by the next "
     "<begin_statements>, so <end> keeps its meaning as the document terminator "
     "\u2014 it stays the stop token and no inference path changes."),
    ("Random draft order. ", "An ascending-F1 ramp lets position alone encode "
     "quality, so \u201clater = better\u201d is learnable without reading the drafts."),
    ("Three-way weight profile. ", "header / draft / final \u2014 swept, because loss on "
     "~13%-precision drafts trains the model to emit WRONG contacts."),
], size=12.5, gap=11)
footer(s, "100,000 documents over 50,000 proteins \u00b7 35,133 packed sequences \u00b7 275 steps/epoch \u00b7 5.3\u00d7 the previous signal.")

# ============================================================ 18. v2 results
s = sl()
title(s, "v2 results \u2014 three falsifications", "553 exp89 proteins, paired. Four arms differ ONLY in loss_weights.")
table(s, 0.75, 2.0, 11.85, [
    ["model (header/draft/final)", "K0", "consensus", "cons \u2212 K0", "wins", "K0 vs base"],
    ["base", "0.3355", "0.2020", "\u22120.1335", "9%", "\u2014"],
    ["mdA  (0 / 0 / 1)", "0.1884", "0.1483", "\u22120.0402", "38%", "\u221243.8%"],
    ["mdB  (0.1 / 0 / 1)", "0.1870", "0.1524", "\u22120.0346", "40%", "\u221244.3%"],
    ["mdC  (0.1 / 0.1 / 1)", "0.1879", "0.1549", "\u22120.0331", "42%", "\u221244.0%"],
    ["mdD  (0.1 / 0.3 / 1)", "0.1886", "0.1509", "\u22120.0378", "40%", "\u221243.8%"],
], [3.6, 1.6, 1.9, 1.7, 1.35, 1.7], size=12.5, row_h=0.37, hilite={(1, 1)})
for i2, (n, t) in enumerate([
        (1, "Loss-weight profile does NOTHING \u2014 all four arms within 0.002 on K0 "
            "(mdB\u2212mdA = \u22120.0014 \u00b1 0.0013: wrong direction, ~40\u00d7 too small)"),
        (2, "Scale does NOT fix forgetting \u2014 5.3\u00d7 more docs gave \u221244% vs v1's \u221241%"),
        (3, "The format cost v1's one success \u2014 v1 GAINED from consensus (+0.024, 63%); "
            "every v2 arm loses (\u22120.033 to \u22120.040)")]):
    y = 4.75 + i2 * 0.62
    circle_num(s, 0.85, y, 0.38, n, fill=CORAL)
    box(s, 1.4, y + 0.04, 11.0, 0.5, t, size=13, color=INK)
footer(s, "bpb misled a third time: +5% (0.3915 \u2192 ~0.411) standing in for a \u221244% task collapse.")

# ============================================================ 19. the pattern
s = sl(dark=True)
title(s, "The variable that actually matters", "Everything else was varied and did not move it.", dark=True)
cd = CategoryChartData()
cd.categories = ["LoRA\n900 proteins", "full FT\n18,750 docs", "full FT\n100,000 docs"]
cd.add_series("K0 loss vs base", (-7, -41, -44))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, In(0.9), In(2.15), In(7.0), In(4.0), cd)
ch = gf.chart
chart_style(ch, legend=False)
ser = ch.plots[0].series[0]
for idx, col in enumerate([GREEN, CORAL, CORAL]):
    ser.points[idx].format.fill.solid(); ser.points[idx].format.fill.fore_color.rgb = col
pl = ch.plots[0]
pl.has_data_labels = True
pl.data_labels.number_format = '0"%"'
pl.data_labels.number_format_is_linked = False
pl.data_labels.font.size = Pt(13); pl.data_labels.font.bold = True; pl.data_labels.font.name = BODY
pl.data_labels.font.color.rgb = WHITE
pl.data_labels.position = XL_LABEL_POSITION.INSIDE_END
for ax in (ch.category_axis, ch.value_axis):
    ax.tick_labels.font.color.rgb = C(0xC9, 0xD4, 0xE8)
# with all-negative bars the category labels default to the TOP, on top of the bars
from pptx.enum.chart import XL_TICK_LABEL_POSITION
ch.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
bullets(s, 8.3, 2.3, 4.15, 4.2, [
    ("LoRA vs full fine-tune. ", "That is the only discriminating variable across "
     "corpus size, fold diversity, document format and loss profile."),
    ("A full FT at lr 1e-4 destroys ~42% of the base contact ability regardless.", ""),
    ("So the next run is LoRA (or a much lower LR) \u2014 not more data, not more "
     "weight tuning.", ""),
], size=13, color=C(0xDE, 0xE7, 0xF5), gap=13, lead=AMBER)
card(s, 8.3, 6.0, 4.15, 0.75, fill=C(0x2E, 0x38, 0x74))
box(s, 8.6, 6.2, 3.6, 0.45, "Still unmeasured: does draft t+1 beat draft t?",
    size=12.5, color=AMBER, bold=True)

# ============================================================ 20. 1M decision
s = sl(dark=True)
title(s, "The 1M push — now a harder call",
      "Cost measured on the 10k batch, not estimated. But scale is not the blocker.", dark=True)
for i, (v, l, col) in enumerate([
        ("14,350", "rollouts per GPU-hour", AMBER),
        ("~1,700", "H100-hours for 1M × 24", WHITE),
        ("~26 h", "wall on 64 shards", C(0x9F, 0xC4, 0xE0))]):
    x = 0.9 + i * 3.95
    card(s, x, 2.2, 3.55, 1.85, fill=C(0x2E, 0x38, 0x74))
    box(s, x, 2.45, 3.55, 0.8, v, size=34, bold=True, font=HEAD, color=col, align=PP_ALIGN.CENTER)
    box(s, x, 3.3, 3.55, 0.5, l, size=12.5, color=C(0xC9, 0xD4, 0xE8), align=PP_ALIGN.CENTER)
bullets(s, 0.9, 4.4, 11.5, 2.3, [
    ("~2× the earlier estimate. ", "The top-k fix roughly doubled rollout length — which "
     "is the point (n_pred 201 vs 95), but it doubles the bill."),
    ("Is 1M the right number? ", "250k proteins is ~420 H100-hours (~26h on 16 shards) and "
     "plausibly buys most of the fold diversity for a quarter of the spend."),
    ("Phase 3 changes the question. ", "The kill criterion is met and the binding "
     "constraint is forgetting, not data volume — so spending ~1,700 H100-hours on the "
     "same recipe would buy a better-trained version of a model that loses to its own base."),
], size=14, color=C(0xDE, 0xE7, 0xF5), gap=12, lead=AMBER)

prs.save("exp163_deck.pptx")
print("saved exp163_deck.pptx", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
